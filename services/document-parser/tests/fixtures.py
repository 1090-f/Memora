"""测试 fixtures：最小合法 PDF / DOCX 构造器。"""

from __future__ import annotations

from io import BytesIO


def build_pdf(text: str = "Hello Docling") -> bytes:
    """手工构造最小合法 PDF（含正确 xref 偏移），文本层可被 pdfium 提取。"""
    objects: list[str] = []
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>".decode())
    objects.append(b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>".decode())
    objects.append(
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>".decode()
    )
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
    objects.append(
        f"<< /Length {len(stream)} >>\nstream\n".encode().decode()
        + stream.decode("latin-1")
        + "\nendstream"
    )
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>".decode())

    out = BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets: list[int] = []
    for i, obj in enumerate(objects, start=1):
        offsets.append(out.tell())
        out.write(f"{i} 0 obj\n".encode())
        out.write(obj.encode("latin-1"))
        out.write(b"\nendobj\n")
    xref_pos = out.tell()
    out.write(f"xref\n0 {len(objects) + 1}\n".encode())
    out.write(b"0000000000 65535 f \n")
    for offset in offsets:
        out.write(f"{offset:010d} 00000 n \n".encode())
    out.write(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF\n".encode())
    return out.getvalue()


def build_docx(
    *,
    heading: str = "第一章",
    paragraph: str = "这是正文段落",
    table_rows: int = 2,
    table_cols: int = 2,
) -> bytes:
    """用 python-docx 构造真实 DOCX：标题、段落、表格。

    python-docx 是 docling 的依赖，生成的文件包含 styles.xml 等必要部件，
    docling Word 后端可以正确识别 Heading1 样式与表格结构。
    """
    from docx import Document

    doc = Document()
    doc.add_heading(heading, level=1)
    doc.add_paragraph(paragraph)
    table = doc.add_table(rows=table_rows, cols=table_cols)
    table.style = "Table Grid"
    for r in range(table_rows):
        for c in range(table_cols):
            table.cell(r, c).text = f"R{r}C{c}"
    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def build_fake_pdf() -> bytes:
    """伪造 PDF：扩展名 .pdf 但内容不是 PDF。"""
    return b"this is not a pdf at all, just some text"


def build_fake_docx() -> bytes:
    """伪造 DOCX：扩展名 .docx 但不是合法 ZIP 容器。"""
    return b"PK\x03\x04 not a real zip content"


def build_xlsx(*, sheet_name: str = "Sheet1", rows: int = 2, cols: int = 2) -> bytes:
    """用 openpyxl 构造真实 XLSX：一个带数据的 Sheet。"""
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = sheet_name
    for r in range(1, rows + 1):
        for c in range(1, cols + 1):
            sheet.cell(row=r, column=c, value=f"R{r}C{c}")
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def build_pptx(*, title: str = "测试演示", lines: int = 2) -> bytes:
    """用 python-pptx 构造真实 PPTX：一个标题 + 正文要点。"""
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    for index in range(lines):
        body.text_frame.add_paragraph()
        body.text_frame.paragraphs[index].text = f"要点 {index + 1}"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_png(*, width: int = 8, height: int = 8, color: tuple[int, int, int] = (255, 0, 0)) -> bytes:
    """用 Pillow 构造最小 PNG。"""
    from PIL import Image

    buffer = BytesIO()
    Image.new("RGB", (width, height), color).save(buffer, format="PNG")
    return buffer.getvalue()


def build_jpeg(*, width: int = 8, height: int = 8) -> bytes:
    """用 Pillow 构造最小 JPEG。"""
    from PIL import Image

    buffer = BytesIO()
    Image.new("RGB", (width, height), (0, 0, 255)).save(buffer, format="JPEG")
    return buffer.getvalue()
